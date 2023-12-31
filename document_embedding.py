import openai
import csv
import PyPDF2
import numpy as np
from openai.embeddings_utils import cosine_similarity
from scipy.spatial import distance_matrix
import docx
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import warnings
import os

_context_path = 'docs'
_embeds = "embeds.csv"
_chunks = "chunks.csv"
# models
EMBEDDING_MODEL = "text-embedding-ada-002"
GPT_35_TURBO_MODEL = "gpt-3.5-turbo"

# Takes a path to a PDF and returns the text contents


def read_pdf_file(file_path):
    pdf_reader = PyPDF2.PdfReader(file_path)
    text = ""
    for page_num in range(len(pdf_reader.pages)):
        page_obj = pdf_reader.pages[page_num]
        pages = page_obj.extract_text()
        text += pages
    return text


def read_word_file(file_path):
    doc = docx.Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)


def read_txt_file(file_path):
    with open(file_path, 'r') as file:
        text = file.read()
    return text


def read_ppt_file(file_path):
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    # Split text into words and join with a single space to remove extra whitespace
                    text = ' '.join(paragraph.text.split())
                    # Only append non-empty text
                    if text.strip():
                        full_text.append(text)
    # Join paragraphs with a single newline
    return '\n'.join(full_text)


def read_epub_file(file_path):
    # Filter out the ebooklib warning
    with warnings.catch_warnings():
        warnings.filterwarnings("ignore", category=UserWarning)
        book = epub.read_epub(file_path)
    full_text = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), 'html.parser')
        text = soup.get_text()
        # Split text into words and join with a single space to remove extra whitespace
        cleaned_text = ' '.join(text.split())
        if cleaned_text.strip():
            full_text.append(cleaned_text)
    return '\n'.join(full_text)


def read_html_file(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        html_content = file.read()
    soup = BeautifulSoup(html_content, "html.parser")
    text = soup.get_text()
    # Split text into words and join with a single space to remove extra whitespace
    cleaned_text = ' '.join(text.split())
    return cleaned_text

# Split the input text into smaller chunks of a specified size.


def split_text(text, chunk_size):
    text_chunks = []
    text_length = len(text)
    start = 0
    while start < text_length:
        end = start + chunk_size
        chunk = text[start:end]
        text_chunks.append(chunk)
        start = end
    return text_chunks


def split_text_lists(texts, chunk_size):
    text_chunks = []
    for text in texts:
        text_length = len(text)
        start = 0
        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]
            text_chunks.append(chunk)
            start = end
    return text_chunks


def create_embeddings(text_chunks):
    embeddings = []
    try:
        prepared_chunks = [chunk.replace("\n", " ") for chunk in text_chunks]
        response = openai.Embedding.create(
            input=prepared_chunks, model=EMBEDDING_MODEL)
        if response and "data" in response:
            for data in response["data"]:
                embeddings.append(data["embedding"])
        return embeddings
    except Exception as e:
        print(f"Error creating embeddings: {e}")
        return None


def write_embeddings_to_csv(embeddings, csv_path):
    with open(csv_path, "w", newline="") as csvfile:
        csv_writer = csv.writer(csvfile)
        for embedding in embeddings:
            csv_writer.writerow(embedding)


def read_embeddings_from_csv(csv_path):
    embeddings = []
    with open(csv_path, "r", newline="") as csvfile:
        csv_reader = csv.reader(csvfile)
        for row in csv_reader:
            embedding = [float(value) for value in row]
            embeddings.append(embedding)
    return embeddings


def write_chunks_to_csv(chunks, csv_path):
    with open(csv_path, "w", encoding="utf-8", newline="") as csv_file:
        writer = csv.writer(csv_file)
        writer.writerow(["chunk"])
        for chunk in chunks:
            writer.writerow([chunk])


def read_chunks_from_csv(csv_path):
    chunks = []
    with open(csv_path, "r", encoding="utf-8", newline="") as csv_file:
        reader = csv.reader(csv_file)
        next(reader)  # Skip header row
        for row in reader:
            chunks.append(row[0])
    return chunks


def calculate_centroid(embeddings):
    centroid = np.mean(embeddings, axis=0)
    return centroid


def closest_embeddings_to_centroid(embeddings, centroid, n=3):
    distances = [distance_matrix([embedding], [centroid])[0][0]
                 for embedding in embeddings]
    closest_indices = np.argpartition(distances, range(n))[:n]
    return closest_indices.tolist()


def search_embeddings(query, embeddings, n=3):
    """
    Search for the most similar embeddings to the given query using cosine similarity.

    Args:
        query (str): The input query.
        embeddings (list): A list of embedding vectors.
        n (int): The number of top results to return.

    Returns:
        list: A list of indices of the top N most similar embeddings.
    """
    query_embedding = create_embeddings([query])[0]
    similarities = [cosine_similarity(
        embedding, query_embedding) for embedding in embeddings]
    # Get the indices of the top N most similar embeddings
    top_indices = np.argsort(similarities)[-n:][::-1]
    return top_indices.tolist()


def retrieve_answer(indices, text_chunks, n=3):
    """
    Retrieve the most relevant text from the text chunks using the provided indices.

    Args:
        indices (list): A list of indices of the most similar embeddings.
        text_chunks (list): A list of text chunks.
        n (int): The number of top answers to return.

    Returns:
        list: A list of the top N most relevant answers.
    """
    if n > len(indices):
        n = len(indices)
    answers = [text_chunks[index] for index in indices[:n]]
    return answers


def summarize_text(embeddings, text_chunks, n=3):
    centroid = calculate_centroid(embeddings)
    closest_indices = closest_embeddings_to_centroid(embeddings, centroid, n)
    summary = retrieve_answer(closest_indices, text_chunks, n)
    return summary


def process_pdfs_and_create_csv(pdf_paths, csv_path, chunk_size=1000):
    all_chunks = []
    all_embeddings = []
    for pdf_path in pdf_paths:
        text = read_pdf_file(pdf_path)
        chunks = split_text(text, chunk_size)
        embeddings = create_embeddings(chunks)
        all_chunks.extend(chunks)
        all_embeddings.extend(embeddings)
    write_embeddings_to_csv(all_embeddings, csv_path)
    return csv_path, all_chunks


def process_docs_and_create_csv(

        embeddings_csv_path=f'{_context_path}/{_embeds}',
        chunks_csv_path=f'{_context_path}/{_chunks}',
        chunk_size=1000):

    file_handlers = {
        ".doc": read_word_file,
        ".docx": read_word_file,
        ".ppt": read_ppt_file,
        ".pptx": read_ppt_file,
        ".epub": read_epub_file,
        ".pdf": read_pdf_file,
        ".html": read_html_file
    }
    all_chunks = []
    all_embeddings = []
    for root, _, files in os.walk(_context_path):
        for file in files:
            doc_path = os.path.join(root, file)
            # Check if the file extension is supported
            file_extension = os.path.splitext(doc_path)[1]
            if file_extension in file_handlers:
                # Read the text content using the appropriate helper function
                text = file_handlers[file_extension](doc_path)
                chunks = split_text(text, chunk_size)
                embeddings = create_embeddings(chunks)
                all_chunks.extend(chunks)
                all_embeddings.extend(embeddings)
            else:
                print(f"Skipping unsupported file type: {doc_path}")
        write_embeddings_to_csv(all_embeddings, embeddings_csv_path)
        write_chunks_to_csv(all_chunks, chunks_csv_path)
        return embeddings_csv_path, chunks_csv_path


def check_embeds(folder):
    embeds_path = os.path.join(folder, "embeds.csv")
    return os.path.isfile(embeds_path)


def summary_agent(prompt):
    completion = openai.ChatCompletion.create(
        model=GPT_35_TURBO_MODEL,
        temperature=0.5,
        messages=[
            {"role": "system", "content": "You give a brief summary of given text. \
                The summary should be concise, informative, and accuratly \
                reflect the contents of the given text.\
                reply only with the summary itself."},
            {"role": "user", "content": prompt},
        ]
    )
    reply_content = completion.choices[0].message.content
    return reply_content


def query_agent(prompt):
    completion = openai.ChatCompletion.create(
        model=GPT_35_TURBO_MODEL,
        temperature=0,
        messages=[
            {"role": "system", "content": "You answer a user's question, given some text as context to help \
                answer the question. The user request will be in the form of a list. The first item in the \
                list is the user's question, the other elements in the list will contain text relavent to \
                answering the question. The answer should be concise, informative, and accuratly reflect \
                the contents of the given text. Do not contradict the contents of the given text in your answer"},
            {"role": "user", "content": prompt},
        ]
    )
    reply_content = completion.choices[0].message.content
    return reply_content


def doc_agent(query):
    # folder = str(folder_paths("./docs"))
    folder = os.path.abspath(_context_path)
    # folder = '/Users/n03an/Documents/projects/playground/python/pdf-gpt-langchain/chatservice/docs'
    embeds_csv_path = os.path.join(folder, _embeds)
    chunks_csv_path = os.path.join(folder, _chunks)
    # create embeddings if not present
    if check_embeds(folder) == False:
        process_docs_and_create_csv(folder, embeds_csv_path, chunks_csv_path)

    embeddings = read_embeddings_from_csv(embeds_csv_path)
    chunks = read_chunks_from_csv(chunks_csv_path)

    # search most similar embeddings to the given query
    index = search_embeddings(query, embeddings)
    answer_chunk = " " + str(retrieve_answer(index, chunks))
    query_with_context = str(query) + answer_chunk
    answer = query_agent(query_with_context)
    return answer
    # return [{"role": "user", "content": query_with_context}, {"role": "assistant", "content": answer}]
